import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Creating a presentation object
prs = Presentation()

# Define colors
color_blue = RGBColor(0, 32, 91)
color_white = RGBColor(255, 255, 255)
color_light_gray = RGBColor(240, 240, 240)
color_dark_gray = RGBColor(51, 51, 51)

# Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Lufthansa Developer Center"
subtitle.text = "API Documentation"

# Set background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = color_blue

# Title style
title.text_frame.paragraphs[0].font.size = Inches(1)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = color_white

# Subtitle style
subtitle.text_frame.paragraphs[0].font.size = Inches(0.5)
subtitle.text_frame.paragraphs[0].font.color.rgb = color_white

# Content Slide
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Getting Started with Lufthansa API"

# Set background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = color_white

# Title style
title.text_frame.paragraphs[0].font.size = Inches(0.8)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = color_blue

# Content style
content.text = "1. Obtain an Access Token\n2. Build a Request\n3. Handle the Response"
content.text_frame.paragraphs[0].font.size = Inches(0.5)
content.text_frame.paragraphs[0].font.color.rgb = color_dark_gray

# Adding a shape with color
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.5), Inches(10), Inches(0.5)
)
shape.fill.solid()
shape.fill.fore_color.rgb = color_blue

prs.save("C:/Users/lululu/mai24_bde_airlines/presentation/Lufthansa_API_Template.pptx")
